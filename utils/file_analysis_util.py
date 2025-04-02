# Import required libraries
import pandas as pd  # For data manipulation (not currently used in this code)
from sklearn.feature_extraction.text import TfidfVectorizer  # For text vectorization
from sklearn.metrics.pairwise import cosine_similarity  # For calculating document similarity
from PyPDF2 import PdfReader  # For reading PDF files
from docx import Document  # For reading Word documents
import pptx  # For reading PowerPoint files
import shutil  # For file operations
import os  # For path operations

def extract_text_from_file(file_path):
    """Extract text content from various file formats.
    
    Args:
        file_path (str): Path to the file to extract text from
        
    Returns:
        str: Extracted text content or empty string if extraction fails
    """
    try:
        # Handle PDF files
        if file_path.lower().endswith('.pdf'):
            with open(file_path, 'rb') as f:  # Open in binary mode for PDFs
                pdf = PdfReader(f)
                # Join text from all pages with spaces
                return " ".join(page.extract_text() for page in pdf.pages)
        
        # Handle Word documents
        elif file_path.lower().endswith(('.doc', '.docx')):
            doc = Document(file_path)
            # Join all paragraph texts with spaces
            return " ".join(para.text for para in doc.paragraphs)
        
        # Handle PowerPoint files
        elif file_path.lower().endswith(('.ppt', '.pptx')):
            prs = pptx.Presentation(file_path)
            text = []
            # Extract text from each shape in each slide
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):  # Check if shape contains text
                        text.append(shape.text)
            return " ".join(text)  # Join all text with spaces
        
        # Handle plain text files
        elif file_path.lower().endswith(('.txt', '.md', '.csv', '.json')):
            with open(file_path, 'r', encoding='utf-8') as f:  # UTF-8 encoding
                return f.read()  # Return full content
        
        # Unsupported file types return empty string
        else:
            return ""
    
    # Return empty string if any error occurs during extraction
    except Exception:
        return ""

def group_related_files(folder_path, output_folder="grouped_files", similarity_threshold=0.5):
    """Group files based on textual similarity using cosine similarity of TF-IDF vectors.
    
    Args:
        folder_path (str): Path to folder containing files to analyze
        output_folder (str): Name of folder to store grouped files (default: "grouped_files")
        similarity_threshold (float): Minimum similarity score for grouping (0-1, default: 0.5)
        
    Returns:
        dict: Dictionary containing:
            - status: "success" or "error"
            - message: Result description
            - groups: List of file groups with metadata
            - output_folder: Absolute path to output directory
    """
    try:
        # Get all files in the specified folder (excluding subdirectories)
        files = [f for f in os.listdir(folder_path) 
                if os.path.isfile(os.path.join(folder_path, f))]
        
        # Return error if folder is empty
        if not files:
            return {"status": "error", "message": "No files found in the specified folder"}
        
        # Extract text content from each file
        file_contents = {}
        for file in files:
            file_path = os.path.join(folder_path, file)
            content = extract_text_from_file(file_path)
            # Only store files with non-empty content
            if content.strip():
                file_contents[file] = content
        
        # Return error if no files had readable content
        if not file_contents:
            return {"status": "error", "message": "No readable content found in any files"}
        
        # Convert text to TF-IDF feature vectors
        vectorizer = TfidfVectorizer()
        tfidf_matrix = vectorizer.fit_transform(file_contents.values())
        
        # Calculate pairwise cosine similarity between documents
        similarity_matrix = cosine_similarity(tfidf_matrix)
        
        # Create output folder (won't raise error if exists)
        os.makedirs(output_folder, exist_ok=True)
        
        # Track grouped files and groups
        grouped = set()  # Files that have been grouped
        groups = []      # List to store group information
        
        # Group similar files
        for i, (file1, content1) in enumerate(file_contents.items()):
            # Skip if file already grouped
            if file1 in grouped:
                continue
                
            # Find all similar files above threshold
            similar_files = [file1]  # Start with current file
            for j, (file2, content2) in enumerate(file_contents.items()):
                # Don't compare with self and check similarity threshold
                if i != j and similarity_matrix[i][j] > similarity_threshold:
                    similar_files.append(file2)
            
            # Only create groups with multiple files
            if len(similar_files) > 1:
                group_name = f"group_{len(groups)+1}"  # Sequential group naming
                group_path = os.path.join(output_folder, group_name)
                os.makedirs(group_path, exist_ok=True)
                
                # Copy all similar files to group folder
                for file in similar_files:
                    src = os.path.join(folder_path, file)
                    dst = os.path.join(group_path, file)
                    shutil.copy2(src, dst)  # copy2 preserves metadata
                    grouped.add(file)  # Mark as grouped
                
                # Store group metadata
                groups.append({
                    "group_name": group_name,
                    "files": similar_files,
                    "similarity_score": max(similarity_matrix[i])  # Highest similarity in group
                })
        
        # Handle ungrouped files
        ungrouped_path = os.path.join(output_folder, "ungrouped")
        os.makedirs(ungrouped_path, exist_ok=True)
        
        # Copy files that didn't meet similarity threshold
        for file in files:
            if file not in grouped:
                src = os.path.join(folder_path, file)
                dst = os.path.join(ungrouped_path, file)
                shutil.copy2(src, dst)
        
        # Return success with grouping results
        return {
            "status": "success",
            "message": f"Grouped {len(groups)} sets of related files",
            "groups": groups,
            "output_folder": os.path.abspath(output_folder)  # Return absolute path
        }
    
    # Handle any exceptions during processing
    except Exception as e:
        return {"status": "error", "message": f"Error processing files: {str(e)}"}